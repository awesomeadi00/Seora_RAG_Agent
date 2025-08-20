import os
import tempfile
import hashlib
from typing import List, Dict, Any, Optional
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract
import chromadb
from chromadb.config import Settings
from rank_bm25 import BM25Okapi
from dotenv import load_dotenv
from openai import OpenAI

# Load the env variables
load_dotenv()

# Hashing function for file storage in ChromaDB
def sha1_bytes(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()

class RAGCore:
    """
    This is the main RAG Agent Class which is structured as such: 
    

    This Agent is responsible for: 
        - 
    """
    # Initialization of OpenAI LLM, Embeddings Model, and setting up the local ChromaDB Database
    def __init__(self, chroma_dir: str, collection_name: str, embed_model: str):
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.embed_model = embed_model
        self.chroma_dir = chroma_dir
        self.collection_name = collection_name

        # Create the ChromaDB database locally if not exists
        os.makedirs(chroma_dir, exist_ok=True)
        self.chroma = chromadb.Client(Settings(is_persistent=True, persist_directory=self.chroma_dir))
        self.collection = self.chroma.get_or_create_collection(self.collection_name)

    # Document Parsing Functions --------------------------------------------------------------------------
    # Helper function read PDF files
    def _parse_pdf(self, file_bytes: bytes) -> str:
        # Create a tmp file to manage data checking through PdfReader
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            reader = PdfReader(tmp.name)    
        # Extract texts from each page in the PDF reader and return
        return "\n\n".join([page.extract_text() or "" for page in reader.pages])

    # Helper function to go through PPTX files
    def _parse_pptx(self, file_bytes: bytes) -> str:
        # Create a tmp file to manage data checking through PPTX Presentation
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            pres = Presentation(tmp.name)
        texts = []
        # Go through each slide in the presentation and extract text and notes
        for i, slide in enumerate(pres.slides, start=1):
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_text.append("NOTES: " + slide.notes_slide.notes_text_frame.text)
            # Append all the text into texts
            if slide_text:
                texts.append(f"[Slide {i}] " + "\n".join(slide_text))
        return "\n\n".join(texts)

    # Helper function to go through CSV File
    def _parse_csv(self, file_bytes: bytes, filename: str) -> str:
        # Create a tmp file to manage data through pandas CSV reader
        with tempfile.NamedTemporaryFile(suffix=".csv", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            df = pd.read_csv(tmp.name)
        # Converts CSV stucture into a string format for LLM input
        return f"[Table: {filename}]\n" + df.to_csv(index=False)

    # Helper function to go through DOCX files
    def _parse_docx(self, file_bytes: bytes) -> str:
        # Create a tmp file to manage data through Document docx class
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
        texts = []
        # For each paragraph in the document, we extract and put into a texts list
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                texts.append(paragraph.text)
        return "\n\n".join(texts)

    # Helper function to go through images through OCR Pytesseract
    def _parse_image(self, file_bytes: bytes, filename: str) -> str:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            image = Image.open(tmp.name)
            # Use OCR to extract text from image
            try:
                text = pytesseract.image_to_string(image)
                return f"[Image: {filename}]\n{text}"
            except Exception as e:
                return f"[Image: {filename}] OCR failed: {str(e)}"
    # ----------------------------------------------------------------------------------------------------
    # Function to chunk the text (default chunk size = 1200, overlap about 12.5%)
    def _chunk_text(self, text: str, chunk_size: int = 1200, overlap: int = 150) -> List[str]:
        # Split the text into a list of words/tokens
        words = text.split()
        chunks = []
        i = 0

        # We go through all the words in the list
        # Append chunks based on the size and account for overlap
        while i < len(words):
            chunk = " ".join(words[i : i + chunk_size])
            chunks.append(chunk)
            i += chunk_size - overlap
        return chunks
        
    # Helper function to create embeddings from the chunked text tokens. 
    def _embed(self, texts: List[str]) -> List[List[float]]:
        resp = self.client.embeddings.create(model=self.embed_model, input=texts)
        return [d.embedding for d in resp.data]
    
    # Re-rank candidate documents using BM25 (keyword-based ranking).
    def _bm25_candidates(self, query: str, docs: List[str], k: int = 10) -> List[int]:
        # For each document, we tokenize/split them and initialize the BM25 with the corpus
        tokenized_corpus = [d.split() for d in docs]
        bm25 = BM25Okapi(tokenized_corpus)
        # Compute BM25 scores for all docs given the tokenized query
        scores = bm25.get_scores(query.split())
        # Return the list sorted giving top-k highest scoring documents
        return sorted(range(len(docs)), key=lambda i: scores[i], reverse=True)[:k]

    # Function to embed the query, and so vector semantic search + BM25 lexical re-ranking to retrieve most relevant content from database
    def _retrieve(self, query: str, top_k: int = 4) -> List[Dict[str, Any]]:
        # Embed the query
        q_emb = self._embed([query])[0]
        # Do vector semantic similarity search within the ChromaDB (pulling top_k - depending on slider how much) 
        vec_res = self.collection.query(query_embeddings=[q_emb], n_results=top_k*3, include=["metadatas", "documents"])
        docs, metas = vec_res["documents"][0], vec_res["metadatas"][0]
        # Re-rank retrieved docs with lexical precision using BM25 ranking
        idx = self._bm25_candidates(query, docs, k=top_k)
        # Return the final top-k chunks with text + metadata
        return [{"text": docs[i], "metadata": metas[i]} for i in idx]

    # This function instantiates the LLM chatmodel and configures the sys prompt and settings
    def _llm_answer(self, question: str, contexts: List[Dict[str, Any]]) -> str:
        sys = """You are an expert startup pitch analyst and business consultant with deep expertise in:
        - Business model analysis and validation
        - Market opportunity assessment (TAM/SAM/SOM)
        - Competitive landscape analysis
        - Financial modeling and projections
        - Investment readiness evaluation
        - Go-to-market strategy development
        - Risk assessment and mitigation

        YOUR ROLE: Analyze startup materials and provide strategic business insights that would be valuable to investors, founders, and business stakeholders.

        RESPONSE STRUCTURE:
        1. **Direct Answer**: Address the question using the provided context snippets and web search results when available
        2. **Business Insights**: Extract startup-relevant insights (market size, competitive advantages, business model, etc.)
        3. **Strategic Analysis**: Provide actionable business recommendations when possible
        4. **Citations**: Always cite sources like (filename#chunk) for uploaded documents and web URLs for online sources

        FOCUS AREAS:
        - Market opportunity and size
        - Competitive positioning and differentiation
        - Business model viability and scalability
        - Financial projections and funding needs
        - Go-to-market strategy effectiveness
        - Risk factors and mitigation strategies

        CONTEXT SOURCES:
        - **Uploaded Documents**: Use these as primary sources for startup-specific analysis
        - **Web Search Results**: Supplement with current market data, industry trends, and competitive intelligence when available
        - **Combined Analysis**: Integrate document insights with web research for comprehensive startup analysis

        IMPORTANT: 
        - Prioritize uploaded document context for startup-specific information
        - Use web search results to enhance analysis with current market data and industry insights
        - Always cite your sources clearly (documents vs. web sources)
        - Maintain a professional, investor-ready tone"""
        
        # We format the context with their metadata for citation tracking
        formatted_ctx = []
        for i, c in enumerate(contexts, start=1):
            meta = c.get("metadata", {})
            formatted_ctx.append(f"[{i}] ({meta.get('source','unknown')}#{meta.get('chunk_index',0)})\n{c.get('text','')}")
        
        # Build the user prompt with the context searched and the question
        ctx_text = "\n\n".join(formatted_ctx) if formatted_ctx else "(no context)"
        prompt = f"Context:\n{ctx_text}\n\nQuestion: {question}\n\nProvide a startup-focused analysis with citations:"
        
        # Instantiate the OpenAI LLM model and return the generated LLM output
        model = os.getenv("MODEL_NAME", "gpt-4o-mini")
        resp = self.client.chat.completions.create(model=model, messages=[{"role":"system","content":sys},{"role":"user","content":prompt}], temperature=0.2)
        return resp.choices[0].message.content

    # Estimate how well the generated answer covers / uses the retrieved contexts.
    def _estimate_coverage(self, answer_text: str, contexts_count: int) -> float:
        # Approximate how many citations were included in the answer
        used = answer_text.count("(") + answer_text.count("[")
        base = min(1.0, contexts_count / 6.0)       # Base score: more contexts retrieved → higher base coverage
        bonus = min(0.4, used * 0.02)               # Bonus score: more citations in the answer → higher coverage, capped at +0.4
        return min(1.0, base + bonus)               # Final coverage = base + bonus (capped at 1.0)

    # Main RAG Class functions to be used outisde ------------------------------------------------------------------------------------------
    # Function to reset the chats by deleting the database and re-instantiating it
    def reset(self):
        self.chroma.delete_collection(self.collection_name)
        self.collection = self.chroma.get_or_create_collection(self.collection_name)

    # Function to index the file 
    def index_file(self, file, filename: Optional[str] = None):
        # Read the filename and hash it
        filename = filename or getattr(file, "name", "upload")
        b = file.read()
        digest = sha1_bytes(b)

        # Get the file extension and depending on the file, call the appropriate parser
        ext = filename.split(".")[-1].lower()
        if ext == "pdf":
            text = self._parse_pdf(b)
        elif ext == "pptx":
            text = self._parse_pptx(b)
        elif ext == "csv":
            text = self._parse_csv(b, filename)
        elif ext == "docx":
            text = self._parse_docx(b)
        elif ext in ["png", "jpg", "jpeg", "tiff", "tif"]:
            text = self._parse_image(b, filename)
        # If no appropriate extension, attempts to convert byte data of the file into a unicode string (last resort)
        else:
            text = b.decode("utf-8", errors="ignore")
        
        # Chunk the extracted text from the file and storing metadata of the file
        chunks = self._chunk_text(text)
        metadata = [{"source": filename, "digest": digest, "chunk_index": i} for i, _ in enumerate(chunks)]
        ids = [f"{digest}_{i}" for i in range(len(chunks))]
        
        # Creating embeddings of each chunk and storing them in the ChromaDB collection
        embeddings = self._embed(chunks)
        self.collection.add(ids=ids, documents=chunks, embeddings=embeddings, metadatas=metadata)

    # Main function for the RAG agent to answer a user query
    def answer(self, question: str, top_k: int = 4, override_contexts:  Optional[List[Dict[str, Any]]] = None) -> Dict[str, Any]:
        # If overrided contexts (merged context from documents + search) is provided then use, else we attempt to retrieve from database
        contexts = override_contexts if override_contexts is not None else self._retrieve(question, top_k=top_k)
        # Generate LLM output of the question and given contexts
        answer = self._llm_answer(question, contexts)
        # Estimate the coverage based on answer and amount of context provided
        coverage = self._estimate_coverage(answer, len(contexts))
        sources = list({c.get("metadata",{}).get("source","unknown") for c in contexts})
        return {"answer": answer, "contexts": contexts, "coverage": coverage, "sources": sources}
    # --------------------------------------------------------------------------------------------------------------------------------------